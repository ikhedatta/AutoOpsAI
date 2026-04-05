"""
AutoOps AI — Hackathon Winning Presentation Generator
Creates a polished PowerPoint deck showcasing the platform.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette ──────────────────────────────────────────────
DARK_BG       = RGBColor(0x0F, 0x17, 0x2A)    # Deep navy
ACCENT_BLUE   = RGBColor(0x00, 0x9B, 0xFF)    # Bright blue
ACCENT_GREEN  = RGBColor(0x00, 0xE6, 0x76)    # Neon green
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)    # Warning orange
ACCENT_RED    = RGBColor(0xFF, 0x4D, 0x4D)    # Alert red
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xB0, 0xB8, 0xC8)
MID_GRAY      = RGBColor(0x6B, 0x72, 0x80)
CARD_BG       = RGBColor(0x1A, 0x24, 0x3B)    # Slightly lighter navy


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, border_color=None):
    """Add a rounded rectangle background shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT_BLUE):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        buNone = pPr.findall(qn('a:buNone'))
        for bn in buNone:
            pPr.remove(bn)
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '▸')
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
        srgbClr.set('val', f'{bullet_color.red:02X}{bullet_color.green:02X}{bullet_color.blue:02X}' if hasattr(bullet_color, 'red') else '009BFF')

    return txBox


def add_tag(slide, left, top, text, bg_color=ACCENT_BLUE, text_color=WHITE, width=None):
    """Add a small tag/badge."""
    w = width or Inches(max(1.2, len(text) * 0.12))
    h = Inches(0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    return shape


# ══════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    # Accent line at top
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    # Main title
    add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                 "AutoOps AI", font_size=54, color=WHITE, bold=True)

    # Subtitle
    add_text_box(slide, Inches(1), Inches(3.0), Inches(10), Inches(0.8),
                 "Your Virtual DevOps Engineer That Never Sleeps", font_size=28,
                 color=ACCENT_BLUE, bold=False)

    # Description
    add_text_box(slide, Inches(1), Inches(4.0), Inches(10), Inches(1.2),
                 "AI-powered infrastructure monitoring, diagnosis, and remediation\n"
                 "with human-in-the-loop approval — fully air-gapped, zero cloud dependencies",
                 font_size=16, color=LIGHT_GRAY)

    # Tags
    tags = [("100% Local AI", ACCENT_GREEN), ("Open Source", ACCENT_BLUE),
            ("Air-Gapped", ACCENT_ORANGE), ("Human-in-the-Loop", RGBColor(0x8B, 0x5C, 0xF6))]
    x = Inches(1)
    for text, color in tags:
        add_tag(slide, x, Inches(5.5), text, bg_color=color)
        x += Inches(max(1.4, len(text) * 0.135))

    # Bottom bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.33), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    add_text_box(slide, Inches(1), Inches(7.05), Inches(11), Inches(0.4),
                 "Hackathon 2026  |  Intelligent Operations Platform",
                 font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_problem(prs):
    """Slide 2: The Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
                 "The Problem", font_size=36, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(7), Inches(0.5),
                 "When production breaks at 2 AM...", font_size=18, color=ACCENT_RED, bold=True)

    # Pain points
    pains = [
        ("🔍  Dig through dashboards", "Engineers scramble across Grafana, logs, and alerts to find what broke"),
        ("📖  Cross-reference runbooks", "Manually search past incidents and documentation for the fix"),
        ("🤝  Get approval", "Wake up another engineer to approve changes to production"),
        ("⚙️  Execute manually", "Type commands by hand in a high-stress, error-prone situation"),
        ("✅  Verify the fix", "Run more manual checks hoping nothing else broke"),
    ]

    y = Inches(1.7)
    for title, desc in pains:
        card = add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(0.85), CARD_BG, border_color=RGBColor(0x2A, 0x34, 0x55))
        add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(3.5), Inches(0.4),
                     title, font_size=16, color=WHITE, bold=True)
        add_text_box(slide, Inches(1.0), y + Inches(0.42), Inches(11), Inches(0.4),
                     desc, font_size=13, color=LIGHT_GRAY)
        y += Inches(0.95)

    # Key stat
    add_shape_bg(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.7), RGBColor(0x1A, 0x0A, 0x0A), border_color=ACCENT_RED)
    add_text_box(slide, Inches(1.0), Inches(6.65), Inches(11), Inches(0.55),
                 "⏱  Average MTTR: 45-90 minutes  |  80% of incidents follow known patterns  |  $14.95B market (2025) growing 25.7% CAGR",
                 font_size=14, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)


def slide_solution(prs):
    """Slide 3: Our Solution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 "The Solution: AutoOps AI", font_size=36, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
                 "A virtual DevOps engineer that monitors, diagnoses, explains, and fixes — with you in control",
                 font_size=16, color=ACCENT_BLUE)

    # Pipeline flow
    steps = [
        ("🔎 DETECT", "Real-time monitoring\n15s polling cycle", ACCENT_BLUE),
        ("🧠 DIAGNOSE", "Playbook + LLM\nhybrid reasoning", RGBColor(0x8B, 0x5C, 0xF6)),
        ("📊 CLASSIFY", "Risk assessment\nLOW / MED / HIGH", ACCENT_ORANGE),
        ("✋ APPROVE", "Human-in-the-loop\ntiered autonomy", ACCENT_GREEN),
        ("⚡ EXECUTE", "Automated fix\nvia provider API", ACCENT_BLUE),
        ("✅ VERIFY", "Outcome validation\n+ audit trail", ACCENT_GREEN),
    ]

    x = Inches(0.5)
    card_w = Inches(1.85)
    for title, desc, color in steps:
        card = add_shape_bg(slide, x, Inches(2.0), card_w, Inches(2.2), CARD_BG, border_color=color)
        add_text_box(slide, x + Inches(0.1), Inches(2.15), card_w - Inches(0.2), Inches(0.45),
                     title, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(2.7), card_w - Inches(0.2), Inches(1.2),
                     desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

        # Arrow between cards
        if x < Inches(9):
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x + card_w + Inches(0.05), Inches(2.9), Inches(0.22), Inches(0.25)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MID_GRAY
            arrow.line.fill.background()

        x += card_w + Inches(0.28)

    # Key differentiator
    add_shape_bg(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.4), CARD_BG, border_color=ACCENT_GREEN)
    add_text_box(slide, Inches(1.0), Inches(4.9), Inches(5), Inches(0.4),
                 "What Makes Us Different", font_size=20, color=ACCENT_GREEN, bold=True)

    diffs = [
        "🗣️  Conversational-First — the agent explains in plain English, no dashboard diving",
        "🛡️  Tiered Autonomy — LOW auto-executes, MEDIUM/HIGH require human approval",
        "📚  Playbook + LLM Hybrid — structured knowledge for known issues, AI for novel ones",
        "🔒  100% Air-Gapped — all LLM inference local via Ollama, zero data leaves your network",
        "🔌  Provider-Agnostic — Docker, Kubernetes, ECS, bare-metal via single interface",
    ]
    add_bullet_text(slide, Inches(1.0), Inches(5.35), Inches(11), Inches(2.0),
                    diffs, font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT_GREEN)


def slide_architecture(prs):
    """Slide 4: Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "System Architecture", font_size=36, color=WHITE, bold=True)

    # Component cards
    components = [
        # Row 1
        [("Metrics Collector", "Prometheus + Loki\n15s polling", ACCENT_BLUE, Inches(0.5), Inches(1.3)),
         ("Agent Engine", "Anomaly Detection\nDiagnosis + Risk", RGBColor(0x8B, 0x5C, 0xF6), Inches(3.5), Inches(1.3)),
         ("Knowledge Base", "8 YAML Playbooks\nPattern Matching", ACCENT_ORANGE, Inches(6.5), Inches(1.3)),
         ("LLM Client", "Ollama Local\nqwen3:4b", ACCENT_GREEN, Inches(9.5), Inches(1.3)),],
        # Row 2
        [("Approval Router", "Risk-Based Routing\nDashboard Cards", ACCENT_GREEN, Inches(0.5), Inches(3.4)),
         ("Remediation Executor", "Provider API\nVerify Outcome", ACCENT_RED, Inches(3.5), Inches(3.4)),
         ("Dashboard UI", "Real-time WebSocket\nChat + Timeline", ACCENT_BLUE, Inches(6.5), Inches(3.4)),
         ("Observability Stack", "Prometheus + Loki\nGrafana", ACCENT_ORANGE, Inches(9.5), Inches(3.4)),],
    ]

    for row in components:
        for name, desc, color, x, y in row:
            card = add_shape_bg(slide, x, y, Inches(2.6), Inches(1.7), CARD_BG, border_color=color)
            add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.3), Inches(0.4),
                         name, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, x + Inches(0.15), y + Inches(0.65), Inches(2.3), Inches(0.9),
                         desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Provider abstraction bar
    add_shape_bg(slide, Inches(0.5), Inches(5.5), Inches(11.6), Inches(0.8), RGBColor(0x0A, 0x1A, 0x2F), border_color=ACCENT_BLUE)
    add_text_box(slide, Inches(0.5), Inches(5.55), Inches(11.6), Inches(0.7),
                 "🔌  Infrastructure Provider Interface  →  Docker Compose  |  Kubernetes  |  ECS  |  Bare Metal",
                 font_size=15, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # Tech stack badges
    add_text_box(slide, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
                 "Python 3.12+  •  FastAPI  •  MongoDB + Beanie  •  Ollama (qwen3:4b)  •  Prometheus  •  Loki  •  Grafana  •  Docker  •  WebSocket",
                 font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_tiered_autonomy(prs):
    """Slide 5: Tiered Autonomy — The Secret Sauce."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "Tiered Autonomy — The Trust Model", font_size=36, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.5),
                 "The right level of automation for every situation", font_size=16, color=LIGHT_GRAY)

    tiers = [
        ("LOW RISK", "Cache clear, log rotation, temp cleanup",
         "Auto-Execute + Notify",
         "Agent acts immediately. You get a notification after it's done.",
         ACCENT_GREEN, Inches(1.7)),
        ("MEDIUM RISK", "Service restart, scale up replicas",
         "Approval Card — 5 min timeout",
         "Agent explains the issue and proposed fix. You approve or deny in the dashboard. Defaults to deny on timeout.",
         ACCENT_ORANGE, Inches(3.4)),
        ("HIGH RISK", "Database restart, failover, scale down",
         "Approval Card — No timeout",
         "Agent presents full context with rollback plan. Requires explicit human approval. No auto-execution ever.",
         ACCENT_RED, Inches(5.1)),
    ]

    for title, examples, action, desc, color, y in tiers:
        add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(1.45), CARD_BG, border_color=color)
        add_tag(slide, Inches(1.0), y + Inches(0.15), title, bg_color=color, width=Inches(1.8))
        add_text_box(slide, Inches(3.0), y + Inches(0.12), Inches(4), Inches(0.35),
                     examples, font_size=13, color=LIGHT_GRAY)
        add_text_box(slide, Inches(7.5), y + Inches(0.12), Inches(4.5), Inches(0.35),
                     action, font_size=14, color=color, bold=True)
        add_text_box(slide, Inches(1.0), y + Inches(0.6), Inches(11), Inches(0.75),
                     desc, font_size=13, color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                 "\"AutoOps AI is the smart middle ground — not fully autonomous like Shoreline, not fully manual like PagerDuty\"",
                 font_size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


def slide_demo_scenarios(prs):
    """Slide 6: Live Demo Scenarios."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "Live Demo — 3 Scenarios", font_size=36, color=WHITE, bold=True)

    scenarios = [
        ("Scenario 1: MongoDB Crash", ACCENT_RED, "HIGH",
         [
             "Chaos script kills MongoDB container",
             "Agent detects within 15 seconds — Flask returns 500s",
             "Agent posts approval card: diagnosis + fix + rollback plan",
             "Human approves → Agent restarts MongoDB → Verifies replica set health",
             "Incident resolved in ~47 seconds with full audit trail",
         ]),
        ("Scenario 2: Redis Memory Full", ACCENT_GREEN, "LOW",
         [
             "Redis memory hits 95% capacity",
             "Agent detects via Prometheus metrics",
             "LOW risk — Agent auto-executes cache eviction, no approval needed",
             "Memory drops to 42% — notification sent to dashboard",
             "Zero human intervention required",
         ]),
        ("Scenario 3: CPU Spike", ACCENT_ORANGE, "MEDIUM",
         [
             "Container CPU stuck above 90% for 2+ minutes",
             "Agent diagnoses: possible resource leak or runaway process",
             "MEDIUM risk — Approval card with 5-min timeout",
             "Collects diagnostics, restarts container on approval",
             "Verifies CPU back to normal levels",
         ]),
    ]

    y = Inches(1.2)
    for title, color, risk, steps in scenarios:
        add_shape_bg(slide, Inches(0.5), y, Inches(12.2), Inches(1.7), CARD_BG, border_color=color)
        add_text_box(slide, Inches(0.7), y + Inches(0.1), Inches(5), Inches(0.4),
                     title, font_size=18, color=color, bold=True)
        add_tag(slide, Inches(11.2), y + Inches(0.12), risk, bg_color=color, width=Inches(1.2))

        step_text = "  →  ".join(steps)
        add_text_box(slide, Inches(0.7), y + Inches(0.55), Inches(11.5), Inches(1.0),
                     step_text, font_size=11, color=LIGHT_GRAY)
        y += Inches(1.85)


def slide_competitive(prs):
    """Slide 7: Competitive Landscape."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "Competitive Landscape", font_size=36, color=WHITE, bold=True)

    # Header row
    add_shape_bg(slide, Inches(0.5), Inches(1.2), Inches(12.2), Inches(0.55), RGBColor(0x0A, 0x1A, 0x2F))
    cols = [("Capability", Inches(0.6), Inches(3.5)),
            ("PagerDuty", Inches(4.2), Inches(1.5)),
            ("Dynatrace", Inches(5.8), Inches(1.5)),
            ("StackStorm", Inches(7.4), Inches(1.5)),
            ("Shoreline", Inches(9.0), Inches(1.5)),
            ("AutoOps AI", Inches(10.6), Inches(1.8))]
    for name, x, w in cols:
        color = ACCENT_BLUE if name == "AutoOps AI" else LIGHT_GRAY
        add_text_box(slide, x, Inches(1.25), w, Inches(0.45),
                     name, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # Data rows
    rows = [
        ("AI-Powered Diagnosis", ["❌", "✅", "❌", "❌", "✅"]),
        ("Conversational Interface", ["❌", "❌", "❌", "⚠️", "✅"]),
        ("Human Approval Gates", ["⚠️", "❌", "❌", "❌", "✅"]),
        ("Tiered Risk Autonomy", ["❌", "❌", "❌", "❌", "✅"]),
        ("Air-Gapped / Local LLM", ["N/A", "❌", "N/A", "❌", "✅"]),
        ("Playbook + LLM Hybrid", ["❌", "❌", "⚠️", "❌", "✅"]),
        ("Provider-Agnostic", ["✅", "✅", "✅", "⚠️", "✅"]),
        ("Open Source", ["❌", "❌", "✅", "❌", "✅"]),
    ]

    y = Inches(1.85)
    for i, (capability, checks) in enumerate(rows):
        bg_color = CARD_BG if i % 2 == 0 else DARK_BG
        add_shape_bg(slide, Inches(0.5), y, Inches(12.2), Inches(0.5), bg_color)
        add_text_box(slide, Inches(0.6), y + Inches(0.05), Inches(3.5), Inches(0.4),
                     capability, font_size=12, color=LIGHT_GRAY)
        for j, check in enumerate(checks):
            x_pos = [Inches(4.2), Inches(5.8), Inches(7.4), Inches(9.0), Inches(10.6)][j]
            w = Inches(1.8) if j == 4 else Inches(1.5)
            clr = ACCENT_GREEN if check == "✅" else (ACCENT_ORANGE if check == "⚠️" else MID_GRAY)
            add_text_box(slide, x_pos, y + Inches(0.05), w, Inches(0.4),
                         check, font_size=13, color=clr, alignment=PP_ALIGN.CENTER)
        y += Inches(0.52)

    add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5),
                 "AutoOps AI is the only platform combining all six capabilities in a single, air-gapped solution",
                 font_size=15, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)


def slide_tech_stack(prs):
    """Slide 8: Tech Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "Technology Stack", font_size=36, color=WHITE, bold=True)

    categories = [
        ("Core Platform", ACCENT_BLUE, [
            "Python 3.12+ — Async-first, rich DevOps ecosystem",
            "FastAPI + Uvicorn — Async REST + WebSocket, auto OpenAPI docs",
            "MongoDB + Motor + Beanie — Async ODM for incidents, approvals, audit",
            "uv — Fast dependency resolution + virtual env management",
        ]),
        ("AI / LLM", RGBColor(0x8B, 0x5C, 0xF6), [
            "Ollama — Air-gapped local inference, zero data leaves network",
            "qwen3:4b (primary) + mistral:7b (fallback)",
            "JSON structured output for reliable parsing",
            "YAML playbooks — version-controlled remediation runbooks",
        ]),
        ("Observability (PLG Stack)", ACCENT_ORANGE, [
            "Prometheus + cAdvisor — Per-container metrics telemetry",
            "Loki + Promtail — Label-indexed log aggregation",
            "Grafana — Unified metrics + logs UI with auto-provisioned dashboards",
            "Agent self-monitoring via /metrics endpoint",
        ]),
        ("Frontend & Delivery", ACCENT_GREEN, [
            "Dark ops-console dashboard — real-time WebSocket events",
            "Vanilla HTML/CSS/JS SPA — zero build step, served by FastAPI",
            "Docker Compose for infrastructure orchestration",
            "236 unit tests, Ruff lint — zero warnings",
        ]),
    ]

    x = Inches(0.5)
    for title, color, items in categories:
        add_shape_bg(slide, x, Inches(1.2), Inches(2.9), Inches(5.8), CARD_BG, border_color=color)
        add_text_box(slide, x + Inches(0.15), Inches(1.35), Inches(2.6), Inches(0.4),
                     title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_bullet_text(slide, x + Inches(0.15), Inches(1.85), Inches(2.6), Inches(4.8),
                        items, font_size=10, color=LIGHT_GRAY, bullet_color=color)
        x += Inches(3.1)


def slide_demo_flow(prs):
    """Slide 9: Demo Flow (what happens in the live demo)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "How It Works — End to End", font_size=36, color=WHITE, bold=True)

    flow_steps = [
        ("0:00", "Healthy Stack", "Dashboard shows all services green\nNginx + App + MongoDB + Redis running", ACCENT_GREEN),
        ("0:15", "Chaos Injection", "MongoDB container killed via chaos script\nSimulates production failure", ACCENT_RED),
        ("0:30", "Agent Detects", "Anomaly detected in next 15s poll cycle\nFlask app returning 500 errors", ACCENT_ORANGE),
        ("0:45", "AI Diagnosis", "Playbook match + LLM explains in plain English\nRisk classified as HIGH", RGBColor(0x8B, 0x5C, 0xF6)),
        ("1:00", "Approval Card", "Rich card in dashboard chat window\nDiagnosis + fix + rollback plan shown", ACCENT_BLUE),
        ("1:15", "Human Approves", "Engineer clicks Approve button\nFull audit trail captured", ACCENT_GREEN),
        ("1:30", "Auto-Remediate", "Agent restarts MongoDB container\nWaits for health check to pass", ACCENT_BLUE),
        ("1:45", "Verified & Logged", "Stack healthy — incident resolved in <2 min\nAudit log + notification sent", ACCENT_GREEN),
    ]

    left_col = flow_steps[:4]
    right_col = flow_steps[4:]

    for col_idx, col in enumerate([left_col, right_col]):
        base_x = Inches(0.5) if col_idx == 0 else Inches(6.7)
        y = Inches(1.2)
        for time, title, desc, color in col:
            # Time badge
            add_tag(slide, base_x, y + Inches(0.12), time, bg_color=color, width=Inches(0.7))
            # Title
            add_text_box(slide, base_x + Inches(0.85), y + Inches(0.05), Inches(2), Inches(0.35),
                         title, font_size=15, color=color, bold=True)
            # Description
            add_text_box(slide, base_x + Inches(0.85), y + Inches(0.4), Inches(5), Inches(0.8),
                         desc, font_size=11, color=LIGHT_GRAY)
            y += Inches(1.3)


def slide_market(prs):
    """Slide 10: Market Opportunity."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "Market Opportunity", font_size=36, color=WHITE, bold=True)

    # Big numbers
    stats = [
        ("$14.95B", "Market Size (2025)", "Incident Response\n& AIOps Market", ACCENT_BLUE),
        ("25.7%", "CAGR Growth", "Projected to $37.33B\nby 2029", ACCENT_GREEN),
        ("80%", "Repeat Incidents", "Follow known patterns\nthat can be automated", ACCENT_ORANGE),
        ("45-90 min", "Avg MTTR", "Current manual incident\nresolution time", ACCENT_RED),
    ]

    x = Inches(0.5)
    for number, label, desc, color in stats:
        add_shape_bg(slide, x, Inches(1.3), Inches(2.9), Inches(2.2), CARD_BG, border_color=color)
        add_text_box(slide, x, Inches(1.5), Inches(2.9), Inches(0.7),
                     number, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(2.2), Inches(2.9), Inches(0.4),
                     label, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(2.6), Inches(2.9), Inches(0.7),
                     desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        x += Inches(3.1)

    # Target users
    add_text_box(slide, Inches(0.8), Inches(4.0), Inches(5), Inches(0.5),
                 "Target Users", font_size=22, color=WHITE, bold=True)
    targets = [
        "DevOps / SRE teams managing microservices infrastructure",
        "Organizations requiring air-gapped operations (government, finance, healthcare)",
        "Teams drowning in alert fatigue with 100+ daily alerts",
        "Companies seeking to reduce MTTR without sacrificing safety controls",
    ]
    add_bullet_text(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5),
                    targets, font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT_BLUE)


def slide_project_status(prs):
    """Slide 11: What We've Built."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "What We've Built", font_size=36, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.4),
                 "Phase 4 Complete — All Core Components Running", font_size=16, color=ACCENT_GREEN, bold=True)

    items = [
        ("Metrics Collector", "Prometheus + Loki polling", "✅ Running", ACCENT_GREEN),
        ("Anomaly Detection Engine", "Threshold + pattern", "✅ Running", ACCENT_GREEN),
        ("Knowledge Base", "8 YAML playbooks loaded", "✅ Loaded", ACCENT_GREEN),
        ("Ollama LLM Client", "qwen3:4b connected", "✅ Connected", ACCENT_GREEN),
        ("Risk Classification", "LOW / MEDIUM / HIGH", "✅ Active", ACCENT_GREEN),
        ("Approval Router", "Auto / Card / Timeout", "✅ Active", ACCENT_GREEN),
        ("Remediation Executor", "Docker provider API", "✅ Active", ACCENT_GREEN),
        ("MongoDB Incident Store", "Full persistence", "✅ Connected", ACCENT_GREEN),
        ("REST API", "12 endpoints serving", "✅ Serving", ACCENT_GREEN),
        ("WebSocket Events", "Live broadcasting", "✅ Broadcasting", ACCENT_GREEN),
        ("Dashboard UI", "Dark ops-console", "✅ Serving", ACCENT_GREEN),
        ("Test Suite", "236 unit tests", "✅ All Green", ACCENT_GREEN),
    ]

    y = Inches(1.5)
    for i, (name, desc, status, color) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        row_y = y + row * Inches(0.75)
        add_shape_bg(slide, x, row_y, Inches(5.9), Inches(0.6), CARD_BG)
        add_text_box(slide, x + Inches(0.15), row_y + Inches(0.08), Inches(2.5), Inches(0.4),
                     name, font_size=12, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(2.7), row_y + Inches(0.08), Inches(1.8), Inches(0.4),
                     desc, font_size=10, color=LIGHT_GRAY)
        add_text_box(slide, x + Inches(4.6), row_y + Inches(0.08), Inches(1.2), Inches(0.4),
                     status, font_size=10, color=color, bold=True, alignment=PP_ALIGN.RIGHT)

    # Quality badges
    add_shape_bg(slide, Inches(0.5), Inches(6.1), Inches(12.2), Inches(0.6), CARD_BG, border_color=ACCENT_BLUE)
    add_text_box(slide, Inches(0.5), Inches(6.15), Inches(12.2), Inches(0.5),
                 "236 Tests Passing  •  Ruff Lint Zero Warnings  •  Graceful Degradation  •  Full Audit Trail  •  12 REST Endpoints  •  WebSocket Live Events",
                 font_size=13, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


def slide_roadmap(prs):
    """Slide 12: Roadmap / Future Vision."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "Roadmap & Future Vision", font_size=36, color=WHITE, bold=True)

    phases = [
        ("NOW — Phase 4 ✅", ACCENT_GREEN, [
            "Full detection → diagnosis → approval → remediation → verify loop",
            "Docker Compose provider, Ollama LLM, MongoDB persistence",
            "Dark ops-console dashboard with real-time WebSocket events",
        ]),
        ("NEXT — Phase 5-6", ACCENT_BLUE, [
            "Kubernetes provider via same abstraction interface",
            "Learning loop — resolved incidents fed back to playbook",
            "MS Teams Adaptive Cards for enterprise approval workflow",
        ]),
        ("FUTURE — Phase 7+", RGBColor(0x8B, 0x5C, 0xF6), [
            "Predictive analysis — detect failures before they happen",
            "Multi-cluster / multi-cloud management",
            "JWT + RBAC authentication, SSO integration",
            "Incident correlation across services (blast radius analysis)",
        ]),
    ]

    y = Inches(1.2)
    for title, color, items in phases:
        add_shape_bg(slide, Inches(0.5), y, Inches(12.2), Inches(1.7), CARD_BG, border_color=color)
        add_text_box(slide, Inches(0.7), y + Inches(0.1), Inches(5), Inches(0.4),
                     title, font_size=18, color=color, bold=True)
        add_bullet_text(slide, Inches(0.7), y + Inches(0.5), Inches(11.5), Inches(1.2),
                        items, font_size=13, color=LIGHT_GRAY, bullet_color=color)
        y += Inches(1.85)

    # Vision statement
    add_shape_bg(slide, Inches(0.5), Inches(6.8), Inches(12.2), Inches(0.5), RGBColor(0x0A, 0x1A, 0x2F), border_color=ACCENT_GREEN)
    add_text_box(slide, Inches(0.5), Inches(6.82), Inches(12.2), Inches(0.45),
                 "Vision: Every ops team has a virtual DevOps engineer — always watching, always explaining, never acting without permission",
                 font_size=14, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)


def slide_closing(prs):
    """Slide 13: Closing / Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Accent line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 "AutoOps AI", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
                 "Your Virtual DevOps Engineer", font_size=28, color=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # Three pillars recap
    pillars = [
        ("🗣️  Conversational-First", "The agent explains —\nyou understand", ACCENT_BLUE),
        ("🛡️  Tiered Autonomy", "Safe automation with\nhuman control", ACCENT_GREEN),
        ("🔒  Fully Air-Gapped", "Zero data leaves\nyour network", ACCENT_ORANGE),
    ]

    x = Inches(1.5)
    for title, desc, color in pillars:
        add_shape_bg(slide, x, Inches(3.6), Inches(3.2), Inches(1.8), CARD_BG, border_color=color)
        add_text_box(slide, x + Inches(0.2), Inches(3.75), Inches(2.8), Inches(0.5),
                     title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(4.3), Inches(2.8), Inches(0.8),
                     desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        x += Inches(3.6)

    add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
                 "\"Monitors. Diagnoses. Explains. Fixes. Learns. — All with you in control.\"",
                 font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Bottom
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.33), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    add_text_box(slide, Inches(1), Inches(7.05), Inches(11), Inches(0.4),
                 "Thank You  |  Questions?  |  Live Demo Available",
                 font_size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    # Set widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Build all slides
    slide_title(prs)          # 1. Title
    slide_problem(prs)        # 2. The Problem
    slide_solution(prs)       # 3. Our Solution
    slide_architecture(prs)   # 4. Architecture
    slide_tiered_autonomy(prs)# 5. Tiered Autonomy
    slide_demo_scenarios(prs) # 6. Demo Scenarios
    slide_competitive(prs)    # 7. Competitive Landscape
    slide_tech_stack(prs)     # 8. Tech Stack
    slide_demo_flow(prs)      # 9. Demo Flow
    slide_market(prs)         # 10. Market Opportunity
    slide_project_status(prs) # 11. What We've Built
    slide_roadmap(prs)        # 12. Roadmap
    slide_closing(prs)        # 13. Closing

    output = "AutoOpsAI_Hackathon_Presentation.pptx"
    prs.save(output)
    print(f"\n✅ Presentation saved: {output}")
    print(f"   → {len(prs.slides)} slides, 16:9 widescreen format")
    print(f"   → Dark theme, professional design")


if __name__ == "__main__":
    main()
